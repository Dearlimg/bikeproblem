#!/usr/bin/env python3
"""
创建项目演示PPT - 美化版
风格：现代、简洁、大气、蓝色主题
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 蓝色主题色 - 更现代的配色
BLUE_DARK = RGBColor(13, 36, 64)          # 深蓝
BLUE_PRIMARY = RGBColor(25, 113, 194)     # 主蓝
BLUE_LIGHT = RGBColor(79, 172, 254)       # 浅蓝
BLUE_ACCENT = RGBColor(0, 150, 255)       # 强调蓝
BLUE_BG = RGBColor(240, 248, 255)         # 背景蓝
WHITE = RGBColor(255, 255, 255)
GRAY_DARK = RGBColor(64, 64, 64)
GRAY_LIGHT = RGBColor(128, 128, 128)
GREEN = RGBColor(46, 204, 113)
ORANGE = RGBColor(255, 159, 64)

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def create_title_slide(prs):
    """创建标题页 - 美化版"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 渐变背景
    add_gradient_background(slide, BLUE_DARK, BLUE_PRIMARY)
    
    # 添加装饰性形状
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BLUE_ACCENT
    top_bar.line.fill.background()
    
    # 底部装饰条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(10), Inches(0.3))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = BLUE_ACCENT
    bottom_bar.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "共享单车租赁预测系统"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "基于机器学习的智能需求预测分析"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = BLUE_LIGHT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.5))
    info_frame = info_box.text_frame
    info_frame.text = "2024 · 机器学习项目"
    info_para = info_frame.paragraphs[0]
    info_para.font.size = Pt(16)
    info_para.font.color.rgb = BLUE_LIGHT
    info_para.alignment = PP_ALIGN.CENTER

def create_section_header(slide, title):
    """创建章节标题"""
    # 背景条
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = BLUE_PRIMARY
    header_bar.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.LEFT

def create_content_slide(prs, title, content_items):
    """创建内容页 - 美化版"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # 章节标题
    create_section_header(slide, title)
    
    # 内容区域
    y_start = 1.5
    for i, item in enumerate(content_items):
        if item.strip() == "":
            y_start += 0.2
            continue
            
        # 项目符号框
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y_start), Inches(0.15), Inches(0.15))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = BLUE_ACCENT
        bullet.line.fill.background()
        
        # 文本
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_start - 0.05), Inches(8), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        
        # 判断是否是标题（粗体）
        if item.startswith("•") or item.startswith("1.") or item.startswith("2.") or item.startswith("3.") or item.startswith("4."):
            text_para.font.size = Pt(20)
            text_para.font.bold = True
            text_para.font.color.rgb = BLUE_DARK
        else:
            text_para.font.size = Pt(18)
            text_para.font.color.rgb = GRAY_DARK
        
        y_start += 0.6

def create_metrics_slide(prs, title, metrics):
    """创建指标展示页 - 美化版"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_BG
    
    # 章节标题
    create_section_header(slide, title)
    
    # 指标卡片
    card_width = 2.2
    card_height = 1.8
    spacing = 0.3
    x_start = 0.6
    y_start = 2.0
    
    for i, (label, value, desc, color) in enumerate(metrics):
        x = x_start + (i % 3) * (card_width + spacing)
        y = y_start + (i // 3) * (card_height + spacing)
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # 数值
        value_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.3), Inches(card_width - 0.2), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.text = value
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = color
        value_para.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.0), Inches(card_width - 0.2), Inches(0.4))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(16)
        label_para.font.bold = True
        label_para.font.color.rgb = BLUE_DARK
        label_para.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.4), Inches(card_width - 0.2), Inches(0.3))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(11)
        desc_para.font.color.rgb = GRAY_LIGHT
        desc_para.alignment = PP_ALIGN.CENTER

def create_feature_importance_slide(prs):
    """创建特征重要性页 - 美化版"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # 章节标题
    create_section_header(slide, "特征重要性分析")
    
    features = [
        ("hr (小时)", "64.6%", "最重要", BLUE_ACCENT),
        ("temp (温度)", "12.0%", "第二重要", BLUE_PRIMARY),
        ("yr (年份)", "8.7%", "业务增长", BLUE_LIGHT),
        ("workingday", "6.1%", "工作日影响", GRAY_DARK)
    ]
    
    y_start = 2.0
    bar_width = 7.0
    
    for i, (name, pct, desc, color) in enumerate(features):
        y = y_start + i * 1.2
        
        # 特征名称
        name_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(2), Inches(0.4))
        name_frame = name_box.text_frame
        name_frame.text = name
        name_para = name_frame.paragraphs[0]
        name_para.font.size = Pt(18)
        name_para.font.bold = True
        name_para.font.color.rgb = BLUE_DARK
        
        # 百分比条背景
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(y + 0.05), Inches(bar_width), Inches(0.3))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = BLUE_BG
        bar_bg.line.fill.background()
        
        # 百分比条
        pct_value = float(pct.replace("%", ""))
        bar_width_actual = bar_width * (pct_value / 100)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(y + 0.05), Inches(bar_width_actual), Inches(0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # 百分比文字
        pct_box = slide.shapes.add_textbox(Inches(3 + bar_width + 0.2), Inches(y), Inches(1), Inches(0.4))
        pct_frame = pct_box.text_frame
        pct_frame.text = pct
        pct_para = pct_frame.paragraphs[0]
        pct_para.font.size = Pt(20)
        pct_para.font.bold = True
        pct_para.font.color.rgb = color
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(3), Inches(y + 0.4), Inches(bar_width), Inches(0.3))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = GRAY_LIGHT

def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """创建两列对比页 - 美化版"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # 章节标题
    create_section_header(slide, title)
    
    # 左列标题框
    left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.2), Inches(0.6))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = BLUE_PRIMARY
    left_header.line.fill.background()
    
    left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(3.8), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.font.size = Pt(22)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = WHITE
    
    # 左列内容
    y = 2.6
    for item in left_items:
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.12), Inches(0.12))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = BLUE_ACCENT
        bullet.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(0.9), Inches(y - 0.05), Inches(3.6), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(16)
        text_para.font.color.rgb = GRAY_DARK
        y += 0.5
    
    # 右列标题框
    right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.8), Inches(4.2), Inches(0.6))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = GREEN
    right_header.line.fill.background()
    
    right_title_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.9), Inches(3.8), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.font.size = Pt(22)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = WHITE
    
    # 右列内容
    y = 2.6
    for item in right_items:
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(y), Inches(0.12), Inches(0.12))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = GREEN
        bullet.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(5.7), Inches(y - 0.05), Inches(3.6), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(16)
        text_para.font.color.rgb = GRAY_DARK
        y += 0.5

def create_conclusion_slide(prs):
    """创建结论页 - 美化版"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 渐变背景
    add_gradient_background(slide, BLUE_DARK, BLUE_PRIMARY)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "结论与展望"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # 主要内容
    content_items = [
        "✓ 成功建立高精度预测模型（R² = 0.92）",
        "✓ 小时因素是核心驱动因素（64.6%）",
        "✓ 为运营决策提供数据支持",
        "",
        "未来优化方向：",
        "• 超参数调优",
        "• 实时预测系统",
        "• 多城市扩展"
    ]
    
    y = 2.5
    for item in content_items:
        if item.strip() == "":
            y += 0.3
            continue
        
        # 图标或项目符号
        if item.startswith("✓"):
            symbol = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(y), Inches(0.2), Inches(0.2))
            symbol.fill.solid()
            symbol.fill.fore_color.rgb = GREEN
            symbol.line.fill.background()
        elif item.startswith("•"):
            symbol = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(y), Inches(0.15), Inches(0.15))
            symbol.fill.solid()
            symbol.fill.fore_color.rgb = BLUE_LIGHT
            symbol.line.fill.background()
        else:
            symbol = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(y), Inches(0.15), Inches(0.15))
            symbol.fill.solid()
            symbol.fill.fore_color.rgb = BLUE_ACCENT
            symbol.line.fill.background()
        
        # 文本
        text_box = slide.shapes.add_textbox(Inches(1.8), Inches(y - 0.05), Inches(7), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(20)
        text_para.font.color.rgb = WHITE
        y += 0.6

def main():
    """创建PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 标题页
    create_title_slide(prs)
    
    # 2. 项目简介
    create_content_slide(prs, "项目简介", [
        "• 使用机器学习方法预测共享单车租赁数量",
        "• 基于历史数据中的环境因素和季节信息",
        "• 帮助优化车辆调度，预测需求高峰",
        "• 提高运营效率和用户体验"
    ])
    
    # 3. 数据概览
    create_two_column_slide(prs, "数据概览",
        "数据来源",
        ["Capital Bikeshare系统", "华盛顿特区", "2011-2012年数据", "17,379条小时数据"],
        "数据特征",
        ["12个特征维度", "时间、天气、温度", "工作日、节假日", "完整的业务数据"]
    )
    
    # 4. 模型性能
    create_metrics_slide(prs, "模型性能指标", [
        ("R² 分数", "0.9207", "解释92%变异", BLUE_ACCENT),
        ("RMSE", "50.11", "均方根误差", BLUE_PRIMARY),
        ("MAE", "31.10", "平均绝对误差", BLUE_LIGHT),
        ("过拟合", "2.3%", "训练-测试差异", GREEN)
    ])
    
    # 5. 特征重要性
    create_feature_importance_slide(prs)
    
    # 6. 核心发现
    create_content_slide(prs, "核心发现", [
        "1. 小时因素是影响需求的最重要因素（64.6%）",
        "   • 时段因素远超温度、天气等因素",
        "   • 早高峰、晚高峰是需求高峰",
        "",
        "2. 模型预测精度高（R² = 0.92）",
        "   • 能够准确预测中高需求水平",
        "   • 为运营提供可靠支持",
        "",
        "3. 运营策略建议",
        "   • 按时段动态调度车辆",
        "   • 早高峰、晚高峰增加投放",
        "   • 深夜时段减少投放"
    ])
    
    # 7. 模型对比
    create_two_column_slide(prs, "模型对比",
        "线性回归",
        ["R²: 0.3880", "RMSE: 139.21", "MAE: 104.80", "简单但性能较低"],
        "随机森林",
        ["R²: 0.9207", "RMSE: 50.11", "MAE: 31.10", "性能优秀"]
    )
    
    # 8. 技术架构
    create_content_slide(prs, "技术架构", [
        "数据层",
        "  • 数据加载与清洗",
        "  • 特征工程与标准化",
        "",
        "模型层",
        "  • 随机森林回归（1000棵树）",
        "  • 最大深度：10层",
        "  • 并行计算优化",
        "",
        "评估层",
        "  • 多种评估指标（R², RMSE, MAE）",
        "  • 残差分析",
        "  • 特征重要性分析"
    ])
    
    # 9. 结论
    create_conclusion_slide(prs)
    
    # 保存
    output_path = "../doc/共享单车租赁预测系统.pptx"
    prs.save(output_path)
    print(f"✅ 美化版PPT已创建：{output_path}")

if __name__ == "__main__":
    main()

