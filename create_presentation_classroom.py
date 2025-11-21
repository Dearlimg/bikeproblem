#!/usr/bin/env python3
"""
创建课堂报告PPT - 10分钟版本
风格：简洁、重点突出、适合课堂展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 蓝色主题
BLUE_DARK = RGBColor(13, 36, 64)
BLUE_PRIMARY = RGBColor(25, 113, 194)
BLUE_LIGHT = RGBColor(79, 172, 254)
BLUE_ACCENT = RGBColor(0, 150, 255)
BLUE_BG = RGBColor(240, 248, 255)
WHITE = RGBColor(255, 255, 255)
GRAY_DARK = RGBColor(64, 64, 64)
GREEN = RGBColor(46, 204, 113)

def create_title_slide(prs):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "基于随机森林的共享单车\n租赁需求预测研究"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "模式识别课程大作业"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = BLUE_LIGHT
    subtitle_para.alignment = PP_ALIGN.CENTER

def create_simple_slide(prs, title, items):
    """简单内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # 标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # 内容
    y = 1.5
    for item in items:
        if item.strip() == "":
            y += 0.2
            continue
        
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.15), Inches(0.15))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = BLUE_ACCENT
        bullet.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y - 0.05), Inches(8), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(24)
        text_para.font.bold = True if item.startswith("•") else False
        text_para.font.color.rgb = BLUE_DARK if item.startswith("•") else GRAY_DARK
        y += 0.7

def create_metrics_slide(prs):
    """指标展示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_BG
    
    # 标题
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "实验结果"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # 指标卡片
    metrics = [
        ("R² 分数", "0.9207", "解释92%变异"),
        ("RMSE", "50.11", "均方根误差"),
        ("MAE", "31.10", "平均绝对误差")
    ]
    
    x_start = 1.5
    card_width = 2.3
    spacing = 0.4
    y = 2.5
    
    for label, value, desc in metrics:
        x = x_start + (metrics.index((label, value, desc)) * (card_width + spacing))
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BLUE_ACCENT
        card.line.width = Pt(3)
        
        # 数值
        value_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.4), Inches(card_width - 0.2), Inches(0.8))
        value_frame = value_box.text_frame
        value_frame.text = value
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(40)
        value_para.font.bold = True
        value_para.font.color.rgb = BLUE_ACCENT
        value_para.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.3), Inches(card_width - 0.2), Inches(0.4))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(20)
        label_para.font.bold = True
        label_para.font.color.rgb = BLUE_DARK
        label_para.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.7), Inches(card_width - 0.2), Inches(0.3))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = GRAY_DARK
        desc_para.alignment = PP_ALIGN.CENTER

def create_feature_slide(prs):
    """特征重要性页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "特征重要性分析"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    features = [
        ("hr (小时)", 64.6, BLUE_ACCENT),
        ("temp (温度)", 12.0, BLUE_PRIMARY),
        ("yr (年份)", 8.7, BLUE_LIGHT),
        ("workingday", 6.1, GREEN)
    ]
    
    y_start = 2.0
    bar_width = 6.0
    
    for name, pct, color in features:
        y = y_start + features.index((name, pct, color)) * 1.1
        
        # 特征名
        name_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(2), Inches(0.4))
        name_frame = name_box.text_frame
        name_frame.text = name
        name_para = name_frame.paragraphs[0]
        name_para.font.size = Pt(22)
        name_para.font.bold = True
        name_para.font.color.rgb = BLUE_DARK
        
        # 进度条背景
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(y + 0.05), Inches(bar_width), Inches(0.35))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = BLUE_BG
        bar_bg.line.fill.background()
        
        # 进度条
        bar_width_actual = bar_width * (pct / 100)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(y + 0.05), Inches(bar_width_actual), Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # 百分比
        pct_box = slide.shapes.add_textbox(Inches(3.5 + bar_width + 0.2), Inches(y), Inches(1), Inches(0.4))
        pct_frame = pct_box.text_frame
        pct_frame.text = f"{pct}%"
        pct_para = pct_frame.paragraphs[0]
        pct_para.font.size = Pt(24)
        pct_para.font.bold = True
        pct_para.font.color.rgb = color

def main():
    """创建PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 标题页
    create_title_slide(prs)
    
    # 2. 研究背景
    create_simple_slide(prs, "研究背景", [
        "• 共享单车运营面临车辆调度不合理问题",
        "• 准确预测需求有助于优化调度",
        "• 机器学习方法可以有效解决此问题"
    ])
    
    # 3. 算法选择
    create_simple_slide(prs, "算法选择：随机森林", [
        "• 集成学习方法，组合多个决策树",
        "• 能够处理非线性关系",
        "• 可以评估特征重要性",
        "• 预测精度高，泛化能力强"
    ])
    
    # 4. 数据与实验
    create_simple_slide(prs, "数据与实验", [
        "• 数据：Capital Bikeshare系统，17,379条记录",
        "• 特征：12个（时间、天气、温度等）",
        "• 模型：随机森林（1000棵树，深度10）",
        "• 评估：R²、RMSE、MAE等指标"
    ])
    
    # 5. 实验结果
    create_metrics_slide(prs)
    
    # 6. 特征重要性
    create_feature_slide(prs)
    
    # 7. 核心发现
    create_simple_slide(prs, "核心发现", [
        "• 小时因素是影响需求的最重要因素（64.6%）",
        "• 模型预测精度高（R² = 0.92）",
        "• 随机森林远优于线性回归",
        "• 为运营决策提供了数据支持"
    ])
    
    # 8. 总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "总结"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    content = [
        "✓ 成功建立高精度预测模型",
        "✓ 识别关键影响因素",
        "✓ 验证算法有效性",
        "✓ 提供实用分析工具"
    ]
    
    y = 2.5
    for item in content:
        symbol = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(y), Inches(0.2), Inches(0.2))
        symbol.fill.solid()
        symbol.fill.fore_color.rgb = GREEN
        symbol.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(2.5), Inches(y - 0.05), Inches(6), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(28)
        text_para.font.color.rgb = WHITE
        y += 0.8
    
    # 保存
    output_path = "课堂报告PPT.pptx"
    prs.save(output_path)
    print(f"✅ 课堂报告PPT已创建：{output_path}（共8页，适合10分钟报告）")

if __name__ == "__main__":
    main()

